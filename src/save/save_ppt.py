from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor


def make_ppt_page(prs, blank_slide, title, data, header, type="north"):
    slide = prs.slides.add_slide(blank_slide)
    write_title(slide, title)
    # 添加表格:rows行数，cols列数，left和top是在PPT中的位置，width是表的列宽，height是表的行高
    table = slide.shapes.add_table(rows=11,
                                   cols=len(header),
                                   left=Cm(3),
                                   top=Cm(2),
                                   width=Cm(19),
                                   height=Cm(15)
                                   )
    table = table.table
    write_header(table, header)
    write_data(table, data, type)


def write_data(table, data, type):
    if type == "north":
        # 按行写入数据,并且设置格式
        for index, row_data in enumerate(data):
            write_cell(index + 1, 0, table, row_data[1]["SName"])
            write_cell(index + 1, 1, table, row_data[0])
            write_cell(index + 1, 2, table, f"{str(round(row_data[1]['ShareSZ_Chg_One'] / 100000000, ndigits=2))} 亿")
            write_cell(index + 1, 3, table, f"{str(round(row_data[1]['LTZB'] * 100, ndigits=2))} %")
            write_cell(index + 1, 4, table, f"{str(round(row_data[1]['LTZB_One'] * 1000, ndigits=2))} ‰")
    elif type == "south":
        for index, row_data in enumerate(data):
            write_cell(index + 1, 0, table, row_data[1]["SNAME"])
            write_cell(index + 1, 1, table, row_data[0])
            write_cell(index + 1, 2, table, f"{str(round(row_data[1]['SHAREHOLDPRICEONE'] / 100000000, ndigits=2))} 亿")
            write_cell(index + 1, 3, table, str(row_data[1]['SHARESRATE']))


def write_title(slide, title):
    # 在指定位置添加文本框
    textbox = slide.shapes.add_textbox(left=Cm(0.2),
                                       top=Cm(0),
                                       width=Cm(20),
                                       height=Cm(1))
    tf = textbox.text_frame
    # 在文本框中写入文字,文字内容为每页PPT第一行最后一列的数据
    textbox.text = title
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.name = '微软雅黑'
    tf.paragraphs[0].font.bold = True


def write_header(table, header):
    # 写入表头,设置表头的格式
    for i, h in enumerate(header):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.size = Pt(9)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.name = '微软雅黑'
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(49, 134, 155)


def write_cell(row, col, table, cell_text):
    cell = table.cell(row, col)
    cell.text = cell_text
    cell.text_frame.paragraphs[0].font.size = Pt(9)
    cell.text_frame.paragraphs[0].font.name = '微软雅黑'
    cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(240, 240, 240)


def save_to_ppt(output_dir, date_str, north_sorted_dict, south_sorted_dict):
    prs = Presentation()
    blank_slide = prs.slide_layouts[6]
    north_header = ["股票名称", "股票代码", "净买入(元)", "持股占流通股比(%)", "增持占流通股比(‰)"]
    make_ppt_page(prs, blank_slide, title="北向资金买入排行", data=north_sorted_dict[-1:-11:-1], header=north_header)
    make_ppt_page(prs, blank_slide, "北向资金买入排行", north_sorted_dict[-11:-21:-1], header=north_header)
    make_ppt_page(prs, blank_slide, "北向资金卖出排行", north_sorted_dict[0:10], header=north_header)
    make_ppt_page(prs, blank_slide, "北向资金卖出排行", north_sorted_dict[10:20], header=north_header)

    south_header = ["股票简称", "股票代码", "净买入(元)", "持股数量占比(%)"]
    make_ppt_page(prs, blank_slide, title="南向资金买入排行", data=south_sorted_dict[-1:-11:-1], header=south_header,
                  type="south")
    make_ppt_page(prs, blank_slide, "南向资金卖出排行", south_sorted_dict[0:10], header=south_header, type="south")

    prs.save(f'{output_dir}/{date_str}  报告.pptx')
